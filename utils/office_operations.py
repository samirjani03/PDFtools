import os
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from pdf2docx import Converter

def pdf_to_word(pdf_path, docx_path):
    # Convert PDF to Word
    cv = Converter(pdf_path)
    cv.convert(docx_path, start=0, end=None)
    cv.close()

def word_to_pdf(docx_path, pdf_path):
    doc = Document(docx_path)
    c = canvas.Canvas(pdf_path, pagesize=letter)
    width, height = letter
    for para in doc.paragraphs:
        c.drawString(72, height - 72, para.text)
        c.showPage()
    c.save()

def ppt_to_pdf(pptx_path, pdf_path):
    prs = Presentation(pptx_path)
    c = canvas.Canvas(pdf_path, pagesize=letter)
    width, height = letter
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                c.drawString(72, height - 72, shape.text)
        c.showPage()
    c.save()

def excel_to_pdf(xlsx_path, pdf_path):
    wb = load_workbook(xlsx_path)
    c = canvas.Canvas(pdf_path, pagesize=letter)
    width, height = letter
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        for row in ws.iter_rows():
            for cell in row:
                c.drawString(72, height - 72, str(cell.value))
            c.showPage()
    c.save()

# Add more office document operations as needed